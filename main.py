# -*- coding: utf-8 -*-
# JACK 2018

from pptx import Presentation
from pptx.util import Inches,Cm
from pptx.chart.data import ChartData
from pptx.shapes.placeholder import ChartPlaceholder
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

class descToPptx:

    # Specify the base template

    def __init__(self, templateFileName="default_template.pptx"):
        self.templateFileName = templateFileName
        self.prs = Presentation(self.templateFileName)

    # From where we gonna take the description

    def initFromFile(self, filePath):
        pass

    def initFromYaml(self, yaml):
        pass

    def initFromJson(self, json):
        pass

    def initFromDic(self, description):
        self.description = description
        self._scrap()

    # Help documenting the template

    def printLayout(self):
        count = 0
        for slide in self.prs.slide_layouts:
            print("- [%d] %s"%(count,slide.name))
            count += 1
            for pl in slide.placeholders:
               print("\t- [%d] %s - %s"%(pl.placeholder_format.idx, pl.placeholder_format.type, pl.name))

    # Json scrapper

    def _scrap(self):
        if self.description:
            for slideDesc in self.description["slides"]:
                print(slideDesc)
                # Create the slide from the right layout
                layout = self.prs.slide_layouts[slideDesc["layout"]]
                slide = self.prs.slides.add_slide(layout)
                # Iterate through placehorlder an fill them
                for index, content in slideDesc["placeholders"].iteritems():
                    plh = slide.placeholders[int(index)]
                    method = None
                    if type(content) is str:
                        method = self._addText
                    elif type(content) is dict:
                        if content["type"] == "img":
                            method = self._addImage
                        else:
                            method = self._addChart
                    if method:
                        method(plh, content)

    # Placeholder building methods

    def _addText(self, plh, content):
        plh.text = content

    def _addChart(self, plh, content):
        chartType = None

        if content["type"] == "pie":
            chartType = XL_CHART_TYPE.PIE
        if content["type"] == "bar":
            chartType = XL_CHART_TYPE.BAR_STACKED

        if chartType:
            # DATA
            chart_data = ChartData()
            chart_data.categories = content["categories"]
            chart_data.add_series('Series 1', content["data"])

            # UI
            graphic_frame = plh.insert_chart(chartType, chart_data)
            chart = graphic_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

    def _addImage(self, plh, content):
        plh.insert_picture(content["path"])

    # Outputing methods

    def writeToDisk(self):
        outName = self.description["name"]
        self.prs.save(outName)

    def someStreamIOReturningGenertionMethod(self):
        pass

if __name__ == "__main__":

    pptdesc = {
        "name": "out.pptx",\
        "slides": [\
            {\
                "layout":0,\
                "placeholders": {\
                    "0": "Ma super présentation",\
                    "1": "elle est bien quand même"\
                },\
                "freeshapes": {}\
            },{\
                "layout":1,\
                "placeholders": {\
                    "0": "Donnée des sous-vêtement",
                    "13": {\
                        "type":"pie",\
                        "categories": ["slip","boxer","string","rien"],\
                        "data":(12,55,4,1)\
                    }\
                },\
                "freeshapes": {}\
            },{\
                "layout":3,\
                "placeholders": {\
                    "0": "Une image de chat et un graph",
                    "13": {\
                        "type":"img",
                        "path":"cattas.jpg"
                    },
                    "14": {\
                        "type":"bar",\
                        "categories": ["noirs","blancs"],\
                        "data":(30,70)\
                    }\
                },\
                "freeshapes": {}\
            }\
        ]\
    }

    dtp = descToPptx()
    dtp.initFromDic(pptdesc)
    dtp.printLayout()
    dtp.writeToDisk()

